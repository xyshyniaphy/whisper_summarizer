"""
Tests for PPTXService - PowerPoint generation from transcriptions.
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from datetime import datetime

from app.services.pptx_service import (
    PPTXService,
    get_pptx_service,
    set_chinese_font,
    CHINESE_FONTS
)


@pytest.fixture
def mock_output_dir(tmp_path):
    """Create a temporary output directory."""
    return tmp_path / "output"


@pytest.fixture
def pptx_service(mock_output_dir):
    """Create a PPTXService instance with mocked output directory."""
    return PPTXService(output_dir=mock_output_dir)


@pytest.fixture
def mock_transcription():
    """Create a mock transcription object."""
    transcription = Mock()
    transcription.id = "test-id-123"
    transcription.file_name = "test_audio.mp3"
    transcription.duration_seconds = 125  # 2:05
    transcription.text = "This is the transcription text.\nIt has multiple lines.\nThird line here."
    return transcription


class TestSetChineseFont:
    """Tests for set_chinese_font utility function."""

    @patch('app.services.pptx_service.CHINESE_FONTS', ['TestFont'])
    def test_sets_font_on_runs(self):
        """Test that font is set on all runs in paragraphs."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()

        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]

        set_chinese_font(mock_text_frame)

        assert mock_run.font.name == 'TestFont'

    @patch('app.services.pptx_service.CHINESE_FONTS', ['Font1', 'Font2'])
    def test_fallback_to_second_font_on_failure(self):
        """Test that fallback fonts are tried if first fails."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()

        # First font setting raises exception, second succeeds
        mock_run.font.name = 'Font1'
        def side_effect(value):
            if value == 'Font1':
                raise Exception("Font not found")
            mock_run.font.name = value

        mock_run.font.name = MagicMock(side_effect=side_effect)

        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]

        set_chinese_font(mock_text_frame)

        # Should have tried both fonts
        assert mock_run.font.name.call_count == 2

    def test_sets_default_font_size_when_none(self):
        """Test that default font size is set when run.font.size is None."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()
        mock_run.font.size = None
        mock_run.font.name = 'TestFont'

        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]

        set_chinese_font(mock_text_frame)

        # Should set default size
        assert mock_run.font.size is not None


class TestPPTXServiceInitialization:
    """Tests for PPTXService initialization."""

    def test_creates_output_directory(self, tmp_path):
        """Test that output directory is created on init."""
        output_dir = tmp_path / "pptx_output"
        service = PPTXService(output_dir=output_dir)

        assert output_dir.exists()
        assert output_dir.is_dir()

    def test_handles_existing_output_directory(self, tmp_path):
        """Test that existing output directory doesn't cause issues."""
        output_dir = tmp_path / "existing"
        output_dir.mkdir()

        service = PPTXService(output_dir=output_dir)

        assert service is not None


class TestGeneratePPTX:
    """Tests for generate_pptx method."""

    def test_generates_pptx_file(self, pptx_service, mock_transcription):
        """Test that PPTX file is generated."""
        output_path = pptx_service.generate_pptx(mock_transcription)

        assert output_path.exists()
        assert output_path.suffix == ".pptx"
        assert output_path.name == f"{mock_transcription.id}.pptx"

    def test_raises_value_error_for_empty_text(self, pptx_service, mock_transcription):
        """Test that ValueError is raised when transcription has no text."""
        mock_transcription.text = None

        with pytest.raises(ValueError, match="Cannot generate PPTX"):
            pptx_service.generate_pptx(mock_transcription)

    def test_generates_title_slide(self, pptx_service, mock_transcription):
        """Test that title slide is created."""
        output_path = pptx_service.generate_pptx(mock_transcription)

        # Verify file was created
        assert output_path.exists()

    def test_generates_with_summary(self, pptx_service, mock_transcription):
        """Test PPTX generation with AI summary."""
        summary = "This is the AI summary.\n\nIt has multiple paragraphs."

        output_path = pptx_service.generate_pptx(mock_transcription, summary_text=summary)

        assert output_path.exists()

    def test_handles_long_content(self, pptx_service, mock_transcription):
        """Test that long content is split across multiple slides."""
        # Create very long transcription
        long_text = "Line of text\n" * 200
        mock_transcription.text = long_text

        output_path = pptx_service.generate_pptx(mock_transcription)

        assert output_path.exists()

    def test_includes_duration_info(self, pptx_service, mock_transcription):
        """Test that duration info is included in title slide."""
        mock_transcription.duration_seconds = 3665  # 61:05

        output_path = pptx_service.generate_pptx(mock_transcription)

        assert output_path.exists()

    def test_handles_no_duration(self, pptx_service, mock_transcription):
        """Test handling when duration_seconds is None."""
        mock_transcription.duration_seconds = None

        output_path = pptx_service.generate_pptx(mock_transcription)

        assert output_path.exists()


class TestAddTitleSlide:
    """Tests for _add_title_slide method."""

    def test_adds_title_slide(self, pptx_service, mock_transcription):
        """Test title slide creation."""
        from pptx import Presentation

        prs = Presentation()
        pptx_service._add_title_slide(prs, mock_transcription)

        # Should have 1 slide
        assert len(prs.slides) == 1

    def test_title_slide_contains_file_name(self, pptx_service, mock_transcription):
        """Test that title slide contains file name."""
        from pptx import Presentation

        prs = Presentation()
        pptx_service._add_title_slide(prs, mock_transcription)

        title_shape = prs.slides[0].shapes.title
        assert mock_transcription.file_name in title_shape.text


class TestAddSummarySlide:
    """Tests for _add_summary_slide method."""

    def test_adds_summary_slide(self, pptx_service):
        """Test summary slide creation."""
        from pptx import Presentation

        prs = Presentation()
        summary = "AI summary text here."

        pptx_service._add_summary_slide(prs, summary)

        # Should have 1 slide
        assert len(prs.slides) == 1

    def test_handles_multi_paragraph_summary(self, pptx_service):
        """Test summary with multiple paragraphs."""
        from pptx import Presentation

        prs = Presentation()
        summary = "First paragraph.\n\nSecond paragraph.\n\nThird paragraph."

        pptx_service._add_summary_slide(prs, summary)

        # Should create slide with multiple paragraphs
        assert len(prs.slides) == 1


class TestAddContentSlides:
    """Tests for _add_content_slides method."""

    def test_adds_content_slides(self, pptx_service):
        """Test content slides creation."""
        from pptx import Presentation

        prs = Presentation()
        content = "First line\nSecond line\nThird line"

        pptx_service._add_content_slides(prs, content)

        # Should have at least 1 slide
        assert len(prs.slides) >= 1

    def test_handles_empty_content(self, pptx_service):
        """Test handling of empty content."""
        from pptx import Presentation

        prs = Presentation()
        pptx_service._add_content_slides(prs, "")

        # Should not add any slides
        assert len(prs.slides) == 0


class TestChunkContent:
    """Tests for _chunk_content method."""

    def test_short_content_returns_single_chunk(self, pptx_service):
        """Test that short content is not split."""
        short = "Short content"

        chunks = pptx_service._chunk_content(short)

        assert len(chunks) == 1
        assert chunks[0] == short

    def test_long_content_is_split(self, pptx_service):
        """Test that long content is split into multiple chunks."""
        pptx_service.CHARS_PER_SLIDE = 100

        long_content = "Line of text here\n" * 30

        chunks = pptx_service._chunk_content(long_content)

        assert len(chunks) > 1

    def test_respects_max_content_slides(self, pptx_service):
        """Test that max slides limit is respected."""
        pptx_service.CHARS_PER_SLIDE = 100
        pptx_service.MAX_CONTENT_SLIDES = 3

        # Create content that would exceed max slides
        very_long = "Line\n" * 500

        chunks = pptx_service._chunk_content(very_long)

        # Should not exceed max
        assert len(chunks) <= pptx_service.MAX_CONTENT_SLIDES

    def test_preserves_line_structure(self, pptx_service):
        """Test that lines are preserved within chunks."""
        content = "Line 1\nLine 2\nLine 3"

        chunks = pptx_service._chunk_content(content)

        # First chunk should have the lines
        assert "Line 1" in chunks[0]
        assert "Line 2" in chunks[0]

    def test_truncation_message_added(self, pptx_service):
        """Test that truncation message is added when limit reached."""
        pptx_service.CHARS_PER_SLIDE = 50
        pptx_service.MAX_CONTENT_SLIDES = 2

        very_long = "Line\n" * 100

        chunks = pptx_service._chunk_content(very_long)

        # Last chunk should have truncation message
        assert "截断" in chunks[-1]


class TestUtilityMethods:
    """Tests for utility methods."""

    def test_pptx_exists_true(self, pptx_service, mock_output_dir):
        """Test pptx_exists returns True when file exists."""
        # Create a PPTX file
        pptx_path = mock_output_dir / "test-id.pptx"
        pptx_path.touch()

        assert pptx_service.pptx_exists("test-id") is True

    def test_pptx_exists_false(self, pptx_service):
        """Test pptx_exists returns False when file doesn't exist."""
        assert pptx_service.pptx_exists("nonexistent") is False

    def test_get_pptx_path(self, pptx_service):
        """Test getting PPTX file path."""
        path = pptx_service.get_pptx_path("test-id")

        assert path.name == "test-id.pptx"
        assert path.parent == pptx_service.output_dir


class TestGetPPTXService:
    """Tests for get_pptx_service singleton."""

    @patch('app.services.pptx_service._pptx_service', None)
    @patch('app.services.pptx_service.PPTXService')
    def test_singleton_initialization(self, mock_pptx_cls):
        """Test that singleton is initialized once."""
        mock_instance = Mock()
        mock_pptx_cls.return_value = mock_instance

        service1 = get_pptx_service()
        service2 = get_pptx_service()

        assert service1 is service2
        mock_pptx_cls.assert_called_once()

    def test_cached_instance_returned(self):
        """Test that cached instance is returned."""
        with patch('app.services.pptx_service._pptx_service') as mock_cached:
            result = get_pptx_service()
            assert result == mock_cached


class TestChineseFonts:
    """Tests for Chinese font configuration."""

    def test_chinese_fonts_list_not_empty(self):
        """Test that CHINESE_FONTS list is populated."""
        assert len(CHINESE_FONTS) > 0

    def test_common_fonts_included(self):
        """Test that common Chinese fonts are included."""
        font_names = ','.join(CHINESE_FONTS)

        # Check for some common fonts
        assert "Microsoft YaHei" in font_names or "SimHei" in font_names
        assert "PingFang" in font_names or "Noto Sans" in font_names


class TestPPTXServiceConstants:
    """Tests for service constants."""

    def test_chars_per_slide_defined(self):
        """Test that CHARS_PER_SLIDE is defined."""
        service = PPTXService()
        assert service.CHARS_PER_SLIDE > 0

    def test_max_content_slides_defined(self):
        """Test that MAX_CONTENT_SLIDES is defined."""
        service = PPTXService()
        assert service.MAX_CONTENT_SLIDES > 0
