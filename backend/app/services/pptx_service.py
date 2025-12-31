"""
PPTX Generation Service

Creates PowerPoint presentations from transcriptions and summaries.
"""

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.models.transcription import Transcription


# Chinese font names (fallback order)
CHINESE_FONTS = [
    "Microsoft YaHei",      # Windows
    "SimHei",               # Windows
    "PingFang SC",          # macOS
    "Noto Sans CJK SC",     # Linux/open-source
    "WenQuanYi Zen Hei",    # Linux
    "Arial Unicode MS",     # Cross-platform fallback
]

def set_chinese_font(text_frame):
    """
    Set Chinese-compatible font for all paragraphs in a text frame.

    Args:
        text_frame: python-pptx text_frame object
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            # Try each font until one works
            for font_name in CHINESE_FONTS:
                try:
                    run.font.name = font_name
                    break
                except Exception:
                    continue
            # Ensure font size is readable
            if run.font.size is None:
                run.font.size = Pt(18)


class PPTXService:
    """Service for generating PowerPoint presentations from transcriptions."""

    # Characters per slide (rough estimate for content fitting)
    CHARS_PER_SLIDE = 800
    # Maximum slides for content (to prevent huge files)
    MAX_CONTENT_SLIDES = 50

    def __init__(self, output_dir: Path = Path("/app/data/output")):
        """
        Initialize PPTX service.

        Args:
            output_dir: Directory to save generated PPTX files
        """
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_pptx(
        self,
        transcription: Transcription,
        summary_text: str | None = None
    ) -> Path:
        """
        Generate a PowerPoint presentation from transcription data.

        Args:
            transcription: Transcription model instance
            summary_text: Optional AI summary text

        Returns:
            Path to the generated PPTX file

        Raises:
            ValueError: If transcription has no content
        """
        if not transcription.original_text:
            raise ValueError("Cannot generate PPTX: transcription has no content")

        prs = Presentation()

        # Slide 1: Title slide
        self._add_title_slide(prs, transcription)

        # Slide 2: AI Summary (if available)
        if summary_text:
            self._add_summary_slide(prs, summary_text)

        # Content slides: Transcription text
        self._add_content_slides(prs, transcription.original_text)

        # Save the presentation
        output_path = self.output_dir / f"{transcription.id}.pptx"
        prs.save(output_path)

        return output_path

    def _add_title_slide(self, prs: Presentation, transcription: Transcription) -> None:
        """Add title slide with file name and metadata."""
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = transcription.file_name
        set_chinese_font(title.text_frame)

        # Set subtitle with date and duration info
        subtitle = slide.placeholders[1]
        duration_info = ""
        if transcription.duration_seconds:
            minutes = int(transcription.duration_seconds // 60)
            seconds = int(transcription.duration_seconds % 60)
            duration_info = f" | 时长: {minutes}:{seconds:02d}"

        subtitle.text = f"转录文档{duration_info} | 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        set_chinese_font(subtitle.text_frame)

    def _add_summary_slide(self, prs: Presentation, summary_text: str) -> None:
        """Add slide with AI summary."""
        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(bullet_slide_layout)

        # Set title
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "AI 摘要"
        set_chinese_font(title_shape.text_frame)

        # Set content
        body_shape = shapes.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.word_wrap = True

        # Split summary into paragraphs and add as bullets
        paragraphs = summary_text.split('\n\n')
        for i, para_text in enumerate(paragraphs):
            if i == 0:
                text_frame.text = para_text.strip()
            else:
                p = text_frame.add_paragraph()
                p.text = para_text.strip()
                p.level = 0

        set_chinese_font(text_frame)

    def _add_content_slides(self, prs: Presentation, content: str) -> None:
        """
        Add content slides with transcription text.

        Long content is split across multiple slides.
        """
        if not content:
            return

        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content

        # Split content into chunks that fit on slides
        chunks = self._chunk_content(content)

        for i, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(bullet_slide_layout)

            # Set title
            shapes = slide.shapes
            title_shape = shapes.title
            if len(chunks) > 1:
                title_shape.text = f"转录内容 ({i}/{len(chunks)})"
            else:
                title_shape.text = "转录内容"
            set_chinese_font(title_shape.text_frame)

            # Set content
            body_shape = shapes.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.word_wrap = True

            # Split chunk into paragraphs
            paragraphs = chunk.split('\n')
            for j, para_text in enumerate(paragraphs):
                para_text = para_text.strip()
                if not para_text:
                    continue

                if j == 0:
                    text_frame.text = para_text
                else:
                    p = text_frame.add_paragraph()
                    p.text = para_text
                    p.level = 0

            set_chinese_font(text_frame)

    def _chunk_content(self, content: str) -> list[str]:
        """
        Split long content into chunks suitable for slides.

        Args:
            content: Full transcription text

        Returns:
            List of content chunks
        """
        if len(content) <= self.CHARS_PER_SLIDE:
            return [content]

        chunks = []
        current_chunk = []
        current_length = 0
        slide_count = 0

        # Split by lines first to preserve paragraph structure
        lines = content.split('\n')

        for line in lines:
            line_length = len(line) + 1  # +1 for newline

            # Check if we need a new slide
            if (current_length + line_length > self.CHARS_PER_SLIDE and
                    current_chunk):
                chunks.append('\n'.join(current_chunk))
                current_chunk = []
                current_length = 0
                slide_count += 1

                # Prevent excessively large files
                if slide_count >= self.MAX_CONTENT_SLIDES:
                    chunks.append(f"\n[内容过长，已截断。前 {slide_count} 张幻灯片已包含前 "
                                 f"{slide_count * self.CHARS_PER_SLIDE} 字符]")
                    break

            current_chunk.append(line)
            current_length += line_length

        # Add remaining content
        if current_chunk:
            chunks.append('\n'.join(current_chunk))

        return chunks

    def pptx_exists(self, transcription_id: str) -> bool:
        """Check if PPTX file exists for given transcription."""
        pptx_path = self.output_dir / f"{transcription_id}.pptx"
        return pptx_path.exists()

    def get_pptx_path(self, transcription_id: str) -> Path:
        """Get the path to PPTX file for given transcription."""
        return self.output_dir / f"{transcription_id}.pptx"


# Singleton instance
_pptx_service: PPTXService | None = None


def get_pptx_service() -> PPTXService:
    """Get or create the singleton PPTX service instance."""
    global _pptx_service
    if _pptx_service is None:
        _pptx_service = PPTXService()
    return _pptx_service
